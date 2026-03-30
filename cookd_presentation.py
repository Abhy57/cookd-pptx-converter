from pptx import Presentation
from pptx.util import Inches
from bs4 import BeautifulSoup

def html_to_powerpoint(html_content, presentation_path):
    # Initialize the Presentation
    prs = Presentation()
    
    # Parse the HTML content
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Creating slides based on the content
    slides = soup.find_all(['h1', 'h2', 'h3', 'p'])  # Adjust tags as needed for slides
    for slide_content in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank layout for custom content
        
        # Setting the title or body text based on the tag
        if slide_content.name == 'h1':
            title = slide.shapes.title
            title.text = slide_content.get_text()
        elif slide_content.name == 'h2' or slide_content.name == 'h3':
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            textbox.text = slide_content.get_text()
        elif slide_content.name == 'p':
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            textbox.text = slide_content.get_text()
    
    # Save the presentation
    prs.save(presentation_path)

# Example Usage
html_content = """
<h1>Slide 1 Title</h1>
<p>This is the content for slide 1.</p>
<h2>Slide 2 Title</h2>
<p>This is the content for slide 2.</p>
<!-- Add more slides as needed -->
"""
presentation_path = "output.pptx"
html_to_powerpoint(html_content, presentation_path)